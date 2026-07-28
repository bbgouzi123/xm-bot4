import os
import logging
from typing import Optional

logger = logging.getLogger(__name__)

def extract_file_content(file_path: str) -> Optional[str]:
    """
    根据文件后缀名，调用相应的 Python 库解析并提取其文本内容。
    当前支持：.pdf, .docx, .xlsx, .pptx, .txt
    """
    if not file_path or not os.path.exists(file_path):
        logger.warning(f"[DocExtractor] 文件不存在或路径为空: {file_path}")
        return None

    _, ext = os.path.splitext(file_path)
    ext = ext.lower()

    try:
        if ext == '.pdf':
            return _extract_pdf(file_path)
        elif ext == '.docx':
            return _extract_docx(file_path)
        elif ext in ('.xlsx', '.xls'):
            return _extract_xlsx(file_path)
        elif ext == '.pptx':
            return _extract_pptx(file_path)
        elif ext == '.txt':
            return _extract_txt(file_path)
        else:
            logger.info(f"[DocExtractor] 不支持的文件格式，仅返回物理路径: {ext}")
            return None
    except Exception as e:
        logger.error(f"[DocExtractor] 解析文件 {file_path} 异常: {e}", exc_info=True)
        return None


def _extract_pdf(file_path: str) -> str:
    """提取 PDF 内容"""
    import pdfplumber
    text_list = []
    with pdfplumber.open(file_path) as pdf:
        # 为了防止 PDF 文件过大导致 Token 爆炸，限制提取前 30 页
        max_pages = min(len(pdf.pages), 30)
        for i in range(max_pages):
            page_text = pdf.pages[i].extract_text()
            if page_text:
                text_list.append(f"--- Page {i+1} ---\n{page_text}")
    return "\n".join(text_list)


def _extract_docx(file_path: str) -> str:
    """提取 Word 文档内容"""
    import docx
    from docx.text.paragraph import Paragraph
    from docx.table import Table
    
    doc = docx.Document(file_path)
    parts = []
    body = doc.element.body
    for child in body:
        if child.tag.endswith('p'):
            p = Paragraph(child, doc)
            txt = p.text.strip()
            if txt:
                parts.append(txt)
        elif child.tag.endswith('tbl'):
            table = Table(child, doc)
            for row in table.rows:
                cells = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        cells.append(cell_text.replace('\n', ' '))
                if cells:
                    parts.append(" | ".join(cells))
                    
    # 额外提取可能存在的文本框数据
    try:
        txbx_texts = []
        for txbx in doc.element.xpath('//w:txbxContent'):
            for p_elem in txbx.xpath('.//w:p'):
                p = Paragraph(p_elem, doc)
                txt = p.text.strip()
                if txt and txt not in parts:
                    txbx_texts.append(txt)
        if txbx_texts:
            parts.append("\n[文本框数据]:")
            parts.extend(txbx_texts)
    except Exception:
        pass
        
    return "\n".join(parts)


def _extract_xlsx(file_path: str) -> str:
    """提取 Excel 表格内容 (限制每个 Sheet 提取前 50 行以防 Token 溢出)"""
    import openpyxl
    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    sheet_texts = []
    for sheet_name in wb.sheetnames[:3]: # 限制最多解析前 3 个 sheet
        sheet = wb[sheet_name]
        row_list = []
        for r_idx, row in enumerate(sheet.iter_rows(values_only=True)):
            if r_idx >= 50: # 限制前 50 行
                row_list.append("... (更多数据已省略)")
                break
            cells = [str(cell).strip() if cell is not None else "" for cell in row]
            if any(cells): # 过滤整行空行
                row_list.append(" | ".join(cells))
        if row_list:
            sheet_texts.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(row_list))
    return "\n\n".join(sheet_texts)


def _extract_text_from_shape(shape, slide_texts):
    """递归提取单个形状（包括表格、组合形状）中的所有文本"""
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            txt = para.text.strip()
            if txt:
                slide_texts.append(txt)
    if shape.has_table:
        for row in shape.table.rows:
            row_texts = []
            for cell in row.cells:
                cell_text = []
                if cell.text_frame:
                    for para in cell.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            cell_text.append(t)
                cell_val = "\n".join(cell_text).strip()
                if cell_val:
                    row_texts.append(cell_val)
            if row_texts:
                slide_texts.append(" | ".join(row_texts))
    if hasattr(shape, "shapes"):
        for sub_shape in shape.shapes:
            _extract_text_from_shape(sub_shape, slide_texts)


def _extract_pptx(file_path: str) -> str:
    """提取 PPT 内容"""
    from pptx import Presentation
    prs = Presentation(file_path)
    slide_texts = []
    for idx, slide in enumerate(prs.slides[:20]): # 限制前 20 张幻灯片
        texts = []
        for shape in slide.shapes:
            _extract_text_from_shape(shape, texts)
        if texts:
            slide_texts.append(f"--- Slide {idx+1} ---\n" + "\n".join(texts))
    return "\n".join(slide_texts)


def _extract_txt(file_path: str) -> str:
    """提取 TXT 文本"""
    # 自动尝试 UTF-8 和 GBK
    for encoding in ("utf-8", "gbk", "utf-16", "ansi"):
        try:
            with open(file_path, "r", encoding=encoding) as f:
                return f.read(100 * 1024) # 限制最多读取 100 KB
        except UnicodeDecodeError:
            continue
    # 兜底二进制读取，只看 ascii
    with open(file_path, "rb") as f:
        return f.read(50 * 1024).decode("ascii", errors="ignore")
