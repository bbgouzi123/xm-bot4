"""
知识库媒体与复杂格式解析模块（PDF, Excel, PowerPoint）
"""
import io
import logging

logger = logging.getLogger(__name__)

_ocr_engine = None


def _get_local_ocr():
    global _ocr_engine
    if _ocr_engine is None:
        try:
            from rapidocr_onnxruntime import RapidOCR
            _ocr_engine = RapidOCR()
        except Exception as e:
            logger.warning(f"[知识库] 初始化本地 RapidOCR 失败: {e}")
    return _ocr_engine


async def _parse_pdf(raw: bytes) -> str:
    """PDF — pdfplumber 优先，PyPDF2/pypdf 后备，若为扫描件则自动触发本地 OCR 提取文字"""
    text = ""
    try:
        import pdfplumber
        with pdfplumber.open(io.BytesIO(raw)) as pdf:
            pages = [p.extract_text() for p in pdf.pages if p.extract_text()]
            if pages:
                text = "\n\n".join(pages)
    except Exception as e:
        logger.warning(f"[知识库] pdfplumber 解析失败: {e}")

    if not text.strip():
        try:
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(raw))
            pages = [p.extract_text() for p in reader.pages if p.extract_text()]
            if pages:
                text = "\n\n".join(pages)
        except Exception as e:
            logger.warning(f"[知识库] pypdf 解析失败: {e}")

    # 如果解析出的文字极其少，通常是扫描件（纯图片 PDF）
    if len(text.strip()) < 50:
        logger.info("[知识库] PDF 解析出的文本内容过少，判定为扫描件或纯图片 PDF，正在尝试触发本地 OCR 提取文字...")
        ocr_engine = _get_local_ocr()
        if ocr_engine:
            ocr_results = []
            
            # 1. 优先尝试使用 pdfplumber 渲染整页为图片以进行高保真 OCR 识别
            try:
                import pdfplumber
                with pdfplumber.open(io.BytesIO(raw)) as ocr_pdf:
                    max_pages = min(len(ocr_pdf.pages), 15)
                    for page_idx in range(max_pages):
                        page_obj = ocr_pdf.pages[page_idx]
                        page_ocr_texts = []
                        try:
                            # 渲染页面为 150 dpi 的 PIL Image
                            page_img = page_obj.to_image(resolution=150).original
                            img_byte_arr = io.BytesIO()
                            page_img.save(img_byte_arr, format='PNG')
                            img_bytes = img_byte_arr.getvalue()
                            
                            res, elapse = ocr_engine(img_bytes)
                            if res:
                                for item in res:
                                    text_line = item[1].strip()
                                    if text_line:
                                        page_ocr_texts.append(text_line)
                        except Exception as page_ex:
                            logger.error(f"[知识库] PDF 页 {page_idx+1} 渲染 OCR 识别失败: {page_ex}")
                        
                        if page_ocr_texts:
                            ocr_results.append(f"【PDF 第 {page_idx+1} 页 本地 OCR 解析内容】：\n" + "\n".join(page_ocr_texts))
            except Exception as plumber_ocr_ex:
                logger.warning(f"[知识库] pdfplumber 页面渲染 OCR 失败: {plumber_ocr_ex}，尝试回退到内嵌图片提取...")

            # 2. 如果 pdfplumber 失败或未识别到内容，使用 pypdf 提取内嵌图片进行 OCR 识别
            if not ocr_results:
                try:
                    from pypdf import PdfReader
                    reader = PdfReader(io.BytesIO(raw))
                    max_pages = min(len(reader.pages), 15)
                    for page_idx in range(max_pages):
                        page = reader.pages[page_idx]
                        images = page.images
                        if not images:
                            continue
                        
                        page_ocr_texts = []
                        for img_idx, img in enumerate(images):
                            img_bytes = img.data
                            try:
                                res, elapse = ocr_engine(img_bytes)
                                if res:
                                    for item in res:
                                        text_line = item[1].strip()
                                        if text_line:
                                            page_ocr_texts.append(text_line)
                            except Exception as ocr_ex:
                                logger.error(f"[知识库] PDF 页 {page_idx+1} 图片 {img_idx+1} 本地 OCR 识别失败: {ocr_ex}")
                        
                        if page_ocr_texts:
                            ocr_results.append(f"【PDF 第 {page_idx+1} 页 图片 OCR 识别内容】：\n" + "\n".join(page_ocr_texts))
                except Exception as pypdf_ex:
                    logger.error(f"[知识库] pypdf 提取图片 OCR 失败: {pypdf_ex}")
            
            if ocr_results:
                text = "\n\n".join(ocr_results)
        else:
            logger.warning("[知识库] 未安装或无法加载本地 RapidOCR，无法对扫描件进行本地 OCR 识别")

    return text


def _parse_xlsx(raw: bytes) -> str:
    """Excel (.xlsx) — 逐 Sheet 逐行提取"""
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        if len(wb.sheetnames) > 1:
            parts.append(f"【{sheet_name}】")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                parts.append(" | ".join(cells))
    wb.close()
    return "\n".join(parts)


def _parse_xls(raw: bytes) -> str:
    """Excel (.xls) 旧格式"""
    import xlrd
    wb = xlrd.open_workbook(file_contents=raw)
    parts = []
    for sheet in wb.sheets():
        if wb.nsheets > 1:
            parts.append(f"【{sheet.name}】")
        for rx in range(sheet.nrows):
            cells = [str(c).strip() for c in sheet.row_values(rx) if str(c).strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_text_from_shape(shape, slide_texts):
    """递归提取单个形状（包括表格、组合形状）中的所有文本"""
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            txt = para.text.strip()
            if txt:
                slide_texts.append(txt)
    if shape.has_table:
        for row in shape.table.rows:
            row_texts = []
            for cell in row.cells:
                cell_text = []
                if cell.text_frame:
                    for para in cell.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            cell_text.append(t)
                cell_val = "\n".join(cell_text).strip()
                if cell_val:
                    row_texts.append(cell_val)
            if row_texts:
                slide_texts.append(" | ".join(row_texts))
    if hasattr(shape, "shapes"):
        for sub_shape in shape.shapes:
            _extract_text_from_shape(sub_shape, slide_texts)


def _parse_pptx(raw: bytes) -> str:
    """PowerPoint (.pptx) — 提取幻灯片中的文本（基于 XML 节点完整提取，避免偏门 Shape 遗漏）"""
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("需要安装 python-pptx: pip install python-pptx")
    
    prs = Presentation(io.BytesIO(raw))
    parts = []
    
    # 绘图空间和展示空间的 XML 命名空间定义
    NS = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'
    }
    
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        
        # 1. 递归提取该 Slide 底层 XML 中的所有文本段落，避免嵌套 Shape 遗漏
        p_elements = slide._element.findall('.//a:p', namespaces=NS)
        for p in p_elements:
            t_elements = p.findall('.//a:t', namespaces=NS)
            p_text = "".join([elem.text for elem in t_elements if elem.text])
            p_text = p_text.strip()
            if p_text:
                slide_texts.append(p_text)
                
        # 2. 如果该页有演讲者备注 (Notes)，也一并提取作为知识库参考
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_texts.append(f"【备注】：\n{notes}")
                
        if slide_texts:
            parts.append(f"【第{i}页】\n" + "\n".join(slide_texts))
            
    return "\n\n".join(parts)
